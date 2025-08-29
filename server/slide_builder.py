# server/slide_builder.py
from pptx import Presentation
from pptx.util import Pt, Inches
from io import BytesIO
from typing import Dict, Any, List

def _collect_template_assets(prs: Presentation):
    theme = {'pictures': []}
    for s in prs.slides:
        for shp in s.shapes:
            try:
                if hasattr(shp, 'image') and getattr(shp, 'image') is not None:
                    theme['pictures'].append(shp.image.blob)
            except Exception:
                pass
    return theme

def _apply_textframe(tf, lines: List[str], font_size_pt: int=20):
    # Clear existing
    try:
        tf.clear()
    except Exception:
        tf.text = ''
    # Fill
    for i, line in enumerate(lines):
        if i==0:
            p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
            p.text = line
        else:
            p = tf.add_paragraph()
            p.text = line
        for r in p.runs:
            r.font.size = Pt(font_size_pt)

def build_presentation_from_outline(template_bytes: bytes, outline: Dict[str, Any], out_stream: BytesIO):
    """
    Build a .pptx in-memory using python-pptx, reusing template masters/layouts/images.
    outline: { slides: [{title, bullets, notes?}], estimated_slide_count }
    """
    base = Presentation(BytesIO(template_bytes))
    theme = _collect_template_assets(base)

    # Start new presentation based on template (preserves theme)
    new = Presentation(BytesIO(template_bytes))

    # choose layouts
    layout_title = new.slide_layouts[0] if len(new.slide_layouts) > 0 else new.slide_layouts[0]
    layout_tc = new.slide_layouts[1] if len(new.slide_layouts) > 1 else layout_title

    slides = outline.get('slides', [])
    for idx, s in enumerate(slides):
        layout = layout_title if idx==0 else layout_tc
        slide = new.slides.add_slide(layout)

        # find title placeholder
        title_tf = None
        body_tf = None
        for shp in slide.shapes:
            try:
                if not shp.is_placeholder:
                    continue
            except Exception:
                continue
            fmt = getattr(shp, 'placeholder_format', None)
            if not fmt:
                continue
            t = fmt.type
            # 1 title, 2 body/content
            if t == 1 and hasattr(shp, 'text_frame'):
                title_tf = shp.text_frame
            if t in (2, 7, 8) and hasattr(shp, 'text_frame'):
                body_tf = shp.text_frame

        if title_tf:
            _apply_textframe(title_tf, [s.get('title', '')], font_size_pt=28)
        if body_tf:
            bullets = s.get('bullets', [])
            if not bullets:
                bullets = ['']
            _apply_textframe(body_tf, bullets, font_size_pt=18)
        else:
            # fallback: add textbox near center
            left = Inches(1); top = Inches(2); width = Inches(8); height = Inches(3)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            _apply_textframe(txBox.text_frame, [s.get('title','')] + s.get('bullets',[]), font_size_pt=18)

        # notes
        try:
            notes_tf = slide.notes_slide.notes_text_frame
            if 'notes' in s and s['notes']:
                notes_tf.text = s['notes']
        except Exception:
            pass

        # try to reuse a picture from template if a picture placeholder exists
        try:
            pic_ph = next((p for p in slide.placeholders if getattr(p, 'placeholder_format', None) and p.placeholder_format.type == 18), None)
            if pic_ph and theme['pictures']:
                blob = theme['pictures'][idx % len(theme['pictures'])]
                slide.shapes.add_picture(BytesIO(blob), pic_ph.left, pic_ph.top, width=pic_ph.width, height=pic_ph.height)
        except Exception:
            pass

    new.save(out_stream)
